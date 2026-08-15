import asyncio
import subprocess
import tempfile
import logging
from typing import Optional
from pathlib import Path
import pdfplumber
import pandas as pd
from io import BytesIO
from pdf2docx import Converter as Pdf2DocxConverter
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

from docx.shared import Inches, Pt
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import re

# ─── CJK font support ───────────────────────────────────────────────────────
# Fonts that render Chinese / Japanese / Korean glyphs reliably.
CJK_FONTS = ["Noto Sans CJK SC", "Noto Sans SC", "SimSun", "STSong", "WenQuanYi Micro Hei", "Arial Unicode MS"]
CJK_LANG_PREFIXES = {"zh", "ja", "ko"}

def _is_cjk_lang(target_lang: str) -> bool:
    if not target_lang:
        return False
    return target_lang.lower().split("-")[0] in CJK_LANG_PREFIXES

def _set_run_cjk_font(run, font_name: str):
    """Set ASCII, hAnsi, eastAsia, and cs font on a docx run element."""
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    r_fonts = r_pr.find(qn("w:rFonts"))
    if r_fonts is None:
        r_fonts = OxmlElement("w:rFonts")
        r_pr.insert(0, r_fonts)
    for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
        r_fonts.set(qn(attr), font_name)

def apply_cjk_font_to_doc(doc, font_name: str):
    """Walk every run in the document and apply CJK-compatible font."""
    for para in doc.paragraphs:
        for run in para.runs:
            _set_run_cjk_font(run, font_name)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        _set_run_cjk_font(run, font_name)

def compact_docx_layout(doc):
    """Tighten margins and spacing on ALL documents to avoid whitespace gaps
    without deleting any XML elements to preserve images/shapes."""
    for section in doc.sections:
        section.top_margin = Inches(0.5)
        section.bottom_margin = Inches(0.5)
        section.left_margin = Inches(0.5)
        section.right_margin = Inches(0.5)

    def _tighten(para, space_val=Pt(2)):
        fmt = para.paragraph_format
        fmt.space_before = space_val
        fmt.space_after = space_val
        fmt.line_spacing = 1.05
        for r in para.runs:
            if r.font.size and r.font.size.pt > 6:
                r.font.size = Pt(max(6.0, r.font.size.pt - 1.0))

    for para in doc.paragraphs:
        if not para.text or not para.text.strip():
            para.paragraph_format.space_before = Pt(0)
            para.paragraph_format.space_after = Pt(0)
            para.paragraph_format.line_spacing = 1.0
        else:
            _tighten(para)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if not para.text or not para.text.strip():
                        para.paragraph_format.space_before = Pt(0)
                        para.paragraph_format.space_after = Pt(0)
                        para.paragraph_format.line_spacing = 1.0
                    else:
                        _tighten(para)


PROTECTED_TERMS = [
    "Python", "TypeScript", "JavaScript", "Java", "Django", "Node.js", "Express.js", 
    "JAX-RS", "React.js", "React", "Angular", "MongoDB", "MySQL", "Elasticsearch", 
    "Redis", "Google BigQuery", "BigQuery", "PySpark", "Kafka", "AWS", "GCP", 
    "Jenkins", "GitHub Actions", "GitHub", "GitLab CI/CD", "GitLab", "Amazon S3", 
    "S3", "Grafana", "New Relic", "Kibana", "Pytest", "Unittest", "Chai", "Mocha", 
    "Git", "Docker", "Kubernetes", "REST API", "REST APIs", "RESTful API", "RESTful APIs",
    "OOP", "SOLID", "AI Harness", "RAG", "VectorDB", "Pinecone", "PyTorch", "TensorFlow", 
    "Hugging Face", "LoRA", "QLoRA", "Apache Cordova", "Ionic", "freeCodeCamp", "HTML", "CSS", "C++", "C#",
    "SQL", "NoSQL", "Pandas", "NumPy", "SDE-2", "SDE-1", "SDE", "API", "APIs", "OMS",
    "Paytm", "Park+", "Clickdee", "MediaShare", "BearTracks", "Bear River Associates", "Bear River", "Bank of America", "InvenTree",
    "Gurugram", "Noida", "Delhi", "New Delhi", "India", "Gurgaon"
]

def protect_text(text: str) -> tuple[str, dict[str, str]]:
    placeholders = {}
    counter = 0

    email_pattern = r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+'
    phone_pattern = r'\+\d{1,3}[\s-]?\d{9,12}'
    url_pattern = r'https?://[^\s]+'
    
    sorted_terms = sorted(PROTECTED_TERMS, key=len, reverse=True)
    term_patterns = []
    for term in sorted_terms:
        escaped = re.escape(term)
        pattern = ""
        if term[0].isalnum() or term[0] == '+':
            pattern += r'\b'
        pattern += escaped
        if term[-1].isalnum() or term[-1] == '+':
            pattern += r'\b'
        term_patterns.append(pattern)
        
    master_pattern = '|'.join([email_pattern, phone_pattern, url_pattern] + term_patterns)
    
    def repl(match):
        nonlocal counter
        matched_str = match.group(0)
        ph = f"[PH_{counter}]"
        placeholders[ph] = matched_str
        counter += 1
        return ph

    protected_text = re.sub(master_pattern, repl, text, flags=re.IGNORECASE)
    return protected_text, placeholders

def restore_text(text: str, placeholders: dict[str, str]) -> str:
    for ph, orig in placeholders.items():
        escaped_ph = re.escape(ph)
        pattern = re.compile(escaped_ph, re.IGNORECASE)
        text = pattern.sub(orig, text)
        
        ph_num = ph.replace("[PH_", "").replace("]", "")
        space_patterns = [
            rf'\[\s*PH\s*_\s*{ph_num}\s*\]',
            rf'\[\s*ph\s*_\s*{ph_num}\s*\]',
            rf'\[\s*PH\s+{ph_num}\s*\]',
            rf'\[\s*ph\s+{ph_num}\s*\]',
            rf'\[\s*{ph_num}\s*\]'
        ]
        for sp in space_patterns:
            text = re.sub(sp, orig, text, flags=re.IGNORECASE)
    return text

def localize_post_translation(text: str, target_lang: str) -> str:
    if not text:
        return text
    target_lang_lower = target_lang.lower()
    
    # Language-agnostic formatting cleanups
    # Fix spaces around hyphens/dashes for number/date ranges (e.g. "2025 - 2026" -> "2025-2026")
    text = re.sub(r'(\d+)\s*[-\u2013\u2014]\s*(\d+)', r'\1-\2', text)
    # Fix spaces around slashes (e.g. "SDE-1 / SDE-2" -> "SDE-1/SDE-2")
    text = re.sub(r'(\w+)\s*/\s*(\w+)', r'\1/\2', text)
    # Fix space after decimals (e.g. "5. 36" or "5 . 36" -> "5.36")
    text = re.sub(r'(\d+)\s*\.\s*(\d+)', r'\1.\2', text)
    # Fix percentage spacing (e.g. "96 %" -> "96%")
    text = re.sub(r'(\d+)\s*%', r'\1%', text)
    # Fix spacing inside parentheses (e.g. "( SDE-2 )" -> "(SDE-2)")
    text = re.sub(r'\u0028\s+', r'(', text)
    text = re.sub(r'\s+\u0029', r')', text)

    if 'zh' in target_lang_lower:
        text = re.sub(r'(\d+)\s*年\s*(\d+)\s*月\s*[-\u2013\u2014至]\s*(\d+)\s*年\s*(\d+)\s*月', r'\1年\2月-\3年\4月', text)
        text = re.sub(r'(\d+)\s*年\s*(\d+)\s*月', r'\1年\2月', text)
        text = re.sub(r'\s*([，。：、；？！（）])\s*', r'\1', text)
    return text

def should_skip_translation(contents: bytes, target_lang: Optional[str]) -> bool:
    if not target_lang or target_lang == "none":
        return True
    try:
        import langdetect
        text = ""
        with pdfplumber.open(BytesIO(contents)) as pdf:
            for page in pdf.pages[:3]:
                text += page.extract_text() or ""
        if not text.strip():
            logger.info("should_skip_translation: No text extracted. Proceeding with translation attempt.")
            return False
        detected = langdetect.detect(text)
        target_prefix = target_lang.lower().split('-')[0]
        detected_prefix = detected.lower().split('-')[0]
        logger.info(f"should_skip_translation: Detected lang '{detected}' (prefix '{detected_prefix}'), target '{target_lang}' (prefix '{target_prefix}')")
        if target_prefix == detected_prefix:
            return True
    except Exception as e:
        logger.warning(f"Failed to auto-detect PDF language: {e}")
    return False

def should_skip_translation_docx(contents: bytes, target_lang: Optional[str]) -> bool:
    if not target_lang or target_lang == "none":
        return True
    try:
        import langdetect
        doc = Document(BytesIO(contents))
        text = "\n".join(p.text for p in doc.paragraphs[:15])
        if not text.strip():
            logger.info("should_skip_translation_docx: No text extracted. Proceeding with translation attempt.")
            return False
        detected = langdetect.detect(text)
        target_prefix = target_lang.lower().split('-')[0]
        detected_prefix = detected.lower().split('-')[0]
        logger.info(f"should_skip_translation_docx: Detected lang '{detected}' (prefix '{detected_prefix}'), target '{target_lang}' (prefix '{target_prefix}')")
        if target_prefix == detected_prefix:
            return True
    except Exception as e:
        logger.warning(f"Failed to auto-detect DOCX language: {e}")
    return False

def is_resume(text: str, filename: Optional[str] = None) -> bool:
    text_lower = text.lower()
    fn_lower = filename.lower() if filename else ""
    
    if "resume" in fn_lower or "cv" in fn_lower or "curriculum" in fn_lower:
        return True
        
    keywords = ["experience", "education", "skills", "projects", "summary", "employment", "professional"]
    match_count = sum(1 for kw in keywords if kw in text_lower)
    if match_count >= 3:
        return True
    return False

def optimize_docx_layout(doc, filename: Optional[str] = None):
    try:
        text = "\n".join(p.text for p in doc.paragraphs[:20])
        for table in doc.tables[:3]:
            for row in table.rows:
                for cell in row.cells:
                    text += "\n" + cell.text
        if not is_resume(text, filename):
            return
        logger.info("Detected resume layout. Applying margin and spacing optimizations to prevent page gaps.")
        compact_docx_layout(doc)
    except Exception as e:
        logger.warning(f"Failed to optimize docx layout: {e}")


logger = logging.getLogger(__name__)


def safe_translate(translator, text: str) -> str:
    if not text or not text.strip():
        return text
    try:
        translated = translator.translate(text)
        if not translated:
            return text
        if "Error 500" in translated or "Server Error" in translated or "That's an error" in translated:
            logger.warning(f"Google Translate returned error page for text: {text[:100]}")
            return text
        return translated
    except Exception as e:
        logger.error(f"Translation failed: {e}")
        return text


def translate_text_helper(text: str, target_lang: str) -> str:
    if not text.strip():
        return text
    from deep_translator import GoogleTranslator
    try:
        translator = GoogleTranslator(source="auto", target=target_lang)
        if len(text) < 4500:
            return safe_translate(translator, text)
        else:
            chunks = []
            current_chunk = ""
            for paragraph in text.split("\n"):
                if len(current_chunk) + len(paragraph) + 1 < 4500:
                    current_chunk += ("\n" if current_chunk else "") + paragraph
                else:
                    if current_chunk:
                        chunks.append(safe_translate(translator, current_chunk))
                    current_chunk = paragraph
            if current_chunk:
                chunks.append(safe_translate(translator, current_chunk))
            return "\n".join(chunks)
    except Exception as e:
        logger.error(f"Translation helper failed: {e}")
        return text


def _pdf_to_text_sync(contents: bytes, target_lang: Optional[str] = None) -> str:
    with pdfplumber.open(BytesIO(contents)) as pdf:
        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    if target_lang and target_lang != "none":
        text = translate_text_helper(text, target_lang)
    return text


async def pdf_to_text(contents: bytes, target_lang: Optional[str] = None) -> str:
    return await asyncio.to_thread(_pdf_to_text_sync, contents, target_lang)


def _pdf_to_csv_sync(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    all_tables = []
    with pdfplumber.open(BytesIO(contents)) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                if table:
                    all_tables.append(pd.DataFrame(table[1:], columns=table[0]))
    if not all_tables:
        raise ValueError("No tables found in PDF")
    combined = pd.concat(all_tables, ignore_index=True)
    if target_lang and target_lang != "none":
        from deep_translator import GoogleTranslator
        translator = GoogleTranslator(source="auto", target=target_lang)
        for col in combined.columns:
            new_col = str(col)
            try:
                if new_col.strip():
                    new_col = safe_translate(translator, new_col)
            except Exception:
                pass
            combined.rename(columns={col: new_col}, inplace=True)
            combined[new_col] = combined[new_col].apply(
                lambda val: safe_translate(translator, str(val)) if val and isinstance(val, str) and val.strip() else val
            )
    return combined.to_csv(index=False).encode("utf-8")


async def pdf_to_csv(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    return await asyncio.to_thread(_pdf_to_csv_sync, contents, target_lang)


def translate_document_paragraphs(doc, translator):
    # Adjust margins of the document to be tighter (0.5 inch) to prevent page gaps/overflow
    for section in doc.sections:
        section.top_margin = Inches(0.5)
        section.bottom_margin = Inches(0.5)
        section.left_margin = Inches(0.5)
        section.right_margin = Inches(0.5)

    def translate_paragraph(paragraph):
        if paragraph.text and paragraph.text.strip():
            orig_text = paragraph.text
            
            # Protect technologies, companies, emails, phone numbers, and URLs
            protected, placeholders = protect_text(orig_text)
            
            # Translate using deep-translator
            translated = safe_translate(translator, protected)
            
            # Restore protected keywords
            translated = restore_text(translated, placeholders)
            
            # Clean up punctuation spacing and localize dates (e.g. for Chinese)
            translated = localize_post_translation(translated, translator.target)
            
            if translated != orig_text:
                # Store original runs formatting details
                run_styles = []
                for r in paragraph.runs:
                    run_styles.append({
                        'bold': r.bold,
                        'italic': r.italic,
                        'underline': r.underline,
                        'font_name': r.font.name,
                        'font_size': r.font.size,
                        'color': r.font.color.rgb if r.font.color else None
                    })
                
                # Clear existing text in all runs
                for r in paragraph.runs:
                    r.text = ""
                
                # Apply translated text to the first run or create a new one
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    first_run.text = translated
                    if run_styles and run_styles[0]['font_size']:
                        # Slightly reduce font size by 1.0pt to ensure translated text fits cleanly
                        first_run.font.size = Pt(max(6.0, run_styles[0]['font_size'].pt - 1.0))
                else:
                    new_run = paragraph.add_run(translated)
                    if run_styles:
                        new_run.bold = run_styles[0]['bold']
                        new_run.italic = run_styles[0]['italic']
                        new_run.underline = run_styles[0]['underline']
                        if run_styles[0]['font_name']:
                            new_run.font.name = run_styles[0]['font_name']
                        if run_styles[0]['font_size']:
                            new_run.font.size = Pt(max(6.0, run_styles[0]['font_size'].pt - 1.0))
                        if run_styles[0]['color']:
                            new_run.font.color.rgb = run_styles[0]['color']

            # Adjust paragraph spacing to be tighter to prevent page gaps
            p_format = paragraph.paragraph_format
            p_format.space_before = Pt(2)
            p_format.space_after = Pt(2)
            p_format.line_spacing = 1.05

    # Process all paragraphs and table cells without deleting empty paragraphs to preserve images/shapes
    for paragraph in doc.paragraphs:
        if paragraph.text and paragraph.text.strip():
            translate_paragraph(paragraph)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.text and paragraph.text.strip():
                        translate_paragraph(paragraph)

    # Always compact layout after translation to prevent whitespace overflow
    compact_docx_layout(doc)

    # Inject CJK-compatible font so Chinese/Japanese/Korean glyphs render properly
    if _is_cjk_lang(translator.target):
        apply_cjk_font_to_doc(doc, CJK_FONTS[0])



def _pdf_to_docx_sync(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf.write(contents)
        pdf_path = tmp_pdf.name
    docx_path = pdf_path.replace(".pdf", ".docx")
    try:
        cv = Pdf2DocxConverter(pdf_path)
        cv.convert(docx_path)
        cv.close()

        doc = Document(docx_path)
        if target_lang and target_lang != "none" and not should_skip_translation(contents, target_lang):
            from deep_translator import GoogleTranslator
            translator = GoogleTranslator(source="auto", target=target_lang)
            translate_document_paragraphs(doc, translator)  # compact + CJK font applied inside
        else:
            # Always compact layout even without translation
            compact_docx_layout(doc)
        doc.save(docx_path)

        with open(docx_path, "rb") as f:
            docx_bytes = f.read()
        return docx_bytes
    finally:
        Path(pdf_path).unlink(missing_ok=True)
        Path(docx_path).unlink(missing_ok=True)



async def pdf_to_docx(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    return await asyncio.to_thread(_pdf_to_docx_sync, contents, target_lang)


def _docx_translate_sync(contents: bytes, target_lang: str) -> bytes:
    if should_skip_translation_docx(contents, target_lang):
        return contents
    doc = Document(BytesIO(contents))
    from deep_translator import GoogleTranslator
    translator = GoogleTranslator(source="auto", target=target_lang)
    # translate_document_paragraphs already calls compact_docx_layout + CJK font internally
    translate_document_paragraphs(doc, translator)
    out_buf = BytesIO()
    doc.save(out_buf)
    return out_buf.getvalue()


def build_pptx_from_docx_text(docx_bytes: bytes, target_lang: Optional[str] = None) -> bytes:
    """Build a PPTX from the text content of a DOCX (one logical group of paragraphs per slide).
    Used for translated presentations where LibreOffice DOCX→PPTX is unreliable."""
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt, Emu
    from pptx.dml.color import RGBColor
    import textwrap

    doc = Document(BytesIO(docx_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Also pull table cell text
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    paragraphs.append(cell_text)

    if not paragraphs:
        paragraphs = ["(No content)"]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # Group paragraphs into slides (~10 lines per slide)
    LINES_PER_SLIDE = 10
    chunks = [paragraphs[i:i+LINES_PER_SLIDE] for i in range(0, len(paragraphs), LINES_PER_SLIDE)]

    for chunk in chunks:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for idx, line in enumerate(chunk):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = PptPt(16)
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            if _is_cjk_lang(target_lang or ""):
                run.font.name = CJK_FONTS[0]
            p.space_after = Emu(100000)

    out_buf = BytesIO()
    prs.save(out_buf)
    return out_buf.getvalue()


async def docx_to_pdf(contents: bytes, target_lang: Optional[str] = None, filename: Optional[str] = None) -> bytes:
    import shutil
    import os
    libreoffice_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice_bin and os.path.exists("/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        libreoffice_bin = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    
    if not libreoffice_bin:
        raise FileNotFoundError("LibreOffice executable not found. Please install LibreOffice (e.g. 'brew install --cask libreoffice').")

    if target_lang and target_lang != "none" and not should_skip_translation_docx(contents, target_lang):
        logger.info(f"docx_to_pdf: translating to '{target_lang}'")
        contents = await asyncio.to_thread(_docx_translate_sync, contents, target_lang)
        logger.info(f"docx_to_pdf: translation complete, new size={len(contents)} bytes")
    elif not target_lang or target_lang == "none" or should_skip_translation_docx(contents, target_lang):
        # Even if not translating, optimize layout if it's a resume to prevent page gaps in LibreOffice
        try:
            doc = Document(BytesIO(contents))
            optimize_docx_layout(doc, filename)
            out_buf = BytesIO()
            doc.save(out_buf)
            contents = out_buf.getvalue()
        except Exception as e:
            logger.warning(f"Failed to apply docx layout optimizations: {e}")

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp_docx:
        tmp_docx.write(contents)
        docx_path = tmp_docx.name
    output_dir = tempfile.mkdtemp()
    try:
        cmd = [
            libreoffice_bin,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            docx_path,
        ]
        proc = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice failed: {stderr.decode()}")
        pdf_file = Path(output_dir) / (Path(docx_path).stem + ".pdf")
        with open(pdf_file, "rb") as f:
            pdf_bytes = f.read()
        return pdf_bytes
    finally:
        Path(docx_path).unlink(missing_ok=True)
        shutil.rmtree(output_dir, ignore_errors=True)


def _txt_to_pdf_sync(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    text = contents.decode("utf-8")
    if target_lang and target_lang != "none":
        text = translate_text_helper(text, target_lang)
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for line in text.split("\n"):
        story.append(Paragraph(line, styles["Normal"]))
        story.append(Spacer(1, 2))
    doc.build(story)
    buffer.seek(0)
    return buffer.read()


async def txt_to_pdf(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    return await asyncio.to_thread(_txt_to_pdf_sync, contents, target_lang)


def _csv_to_pdf_sync(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    df = pd.read_csv(BytesIO(contents))
    if target_lang and target_lang != "none":
        from deep_translator import GoogleTranslator
        translator = GoogleTranslator(source="auto", target=target_lang)
        for col in df.columns:
            new_col = str(col)
            try:
                if new_col.strip():
                    new_col = safe_translate(translator, new_col)
            except Exception:
                pass
            df.rename(columns={col: new_col}, inplace=True)
            df[new_col] = df[new_col].apply(
                lambda val: safe_translate(translator, str(val)) if val and isinstance(val, str) and val.strip() else val
            )
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    elements = []
    data = [df.columns.tolist()] + df.values.tolist()
    table = Table(data)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("BACKGROUND", (0, 1), (-1, -1), colors.beige),
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
            ]
        )
    )
    elements.append(table)
    doc.build(elements)
    buffer.seek(0)
    return buffer.read()


async def csv_to_pdf(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    return await asyncio.to_thread(_csv_to_pdf_sync, contents, target_lang)


def _pdf_to_excel_sync(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    all_tables = []
    with pdfplumber.open(BytesIO(contents)) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                if table:
                    all_tables.append(pd.DataFrame(table[1:], columns=table[0]))

    if not all_tables:
        # Graceful fallback: extract plain text into a single-column sheet
        logger.warning("No tables found in PDF — falling back to text extraction for Excel output.")
        with pdfplumber.open(BytesIO(contents)) as pdf:
            lines = []
            for page in pdf.pages:
                text = page.extract_text() or ""
                lines.extend([l for l in text.splitlines() if l.strip()])
        combined = pd.DataFrame(lines, columns=["Content"])
    else:
        combined = pd.concat(all_tables, ignore_index=True)

    if target_lang and target_lang != "none":
        from deep_translator import GoogleTranslator
        translator = GoogleTranslator(source="auto", target=target_lang)
        new_cols = {}
        for col in combined.columns:
            new_col = str(col)
            try:
                if new_col.strip():
                    new_col = safe_translate(translator, new_col)
            except Exception:
                pass
            new_cols[col] = new_col
        combined.rename(columns=new_cols, inplace=True)
        for col in combined.columns:
            combined[col] = combined[col].apply(
                lambda val: safe_translate(translator, str(val)) if val and isinstance(val, str) and val.strip() else val
            )
    out_buf = BytesIO()
    with pd.ExcelWriter(out_buf, engine='openpyxl') as writer:
        combined.to_excel(writer, index=False)
    return out_buf.getvalue()


async def pdf_to_excel(contents: bytes, target_lang: Optional[str] = None) -> bytes:
    return await asyncio.to_thread(_pdf_to_excel_sync, contents, target_lang)

def _pdf_to_ppt_sync(contents: bytes) -> bytes:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    
    doc = fitz.open(stream=contents, filetype="pdf")
    prs = Presentation()
    
    # Remove default slides
    for i in range(len(prs.slides)-1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
        
    blank_slide_layout = prs.slide_layouts[6]
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=150)
        img_data = pix.tobytes("png")
        
        img_fp = BytesIO(img_data)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        width_in = page.rect.width / 72.0
        height_in = page.rect.height / 72.0
        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)
        
        slide.shapes.add_picture(img_fp, 0, 0, width=Inches(width_in), height=Inches(height_in))
        
    out_buf = BytesIO()
    prs.save(out_buf)
    return out_buf.getvalue()

async def pdf_to_ppt(contents: bytes) -> bytes:
    return await asyncio.to_thread(_pdf_to_ppt_sync, contents)

async def convert_via_libreoffice(contents: bytes, ext: str) -> bytes:
    """Convert any supported office format to PDF using LibreOffice.
    Includes file-existence check and a fallback command if the primary
    conversion exits without producing output."""
    import shutil
    import os
    libreoffice_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice_bin and os.path.exists("/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        libreoffice_bin = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

    if not libreoffice_bin:
        raise FileNotFoundError("LibreOffice executable not found. Please install LibreOffice.")

    suffix = f".{ext.replace('.', '')}"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp_file:
        tmp_file.write(contents)
        input_path = tmp_file.name

    output_dir = tempfile.mkdtemp()
    try:
        cmd = [
            libreoffice_bin, "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            input_path,
        ]
        proc = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=60)

        pdf_file = Path(output_dir) / (Path(input_path).stem + ".pdf")

        if proc.returncode != 0 or not pdf_file.exists():
            logger.warning(
                f"LibreOffice primary ({ext}→pdf) failed (code={proc.returncode}) or produced no output. "
                "Trying fallback without filter flags."
            )
            # Remove any partial output
            pdf_file.unlink(missing_ok=True)
            for f in Path(output_dir).glob("*.pdf"):
                f.unlink(missing_ok=True)

            cmd_fb = [
                libreoffice_bin, "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                input_path,
            ]
            proc2 = await asyncio.create_subprocess_exec(
                *cmd_fb, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )
            _, stderr2 = await asyncio.wait_for(proc2.communicate(), timeout=60)
            if proc2.returncode != 0:
                raise RuntimeError(f"LibreOffice fallback failed: {stderr2.decode()}")

        # Find the output file
        if not pdf_file.exists():
            found = list(Path(output_dir).glob("*.pdf"))
            if found:
                pdf_file = found[0]
            else:
                raise FileNotFoundError(f"Converted PDF not found for input format .{ext}")

        with open(pdf_file, "rb") as f:
            return f.read()
    finally:
        Path(input_path).unlink(missing_ok=True)
        import shutil as _shutil
        _shutil.rmtree(output_dir, ignore_errors=True)

def _epub_to_pdf_sync(contents: bytes) -> bytes:
    import zipfile
    from bs4 import BeautifulSoup
    
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    with zipfile.ZipFile(BytesIO(contents)) as z:
        html_files = [f for f in z.namelist() if f.endswith(('.html', '.xhtml', '.htm'))]
        html_files.sort()
        for hf in html_files:
            try:
                html_data = z.read(hf).decode("utf-8", errors="ignore")
                soup = BeautifulSoup(html_data, "html.parser")
                paragraphs = soup.find_all(['p', 'h1', 'h2', 'h3'])
                for p in paragraphs:
                    text = p.get_text().strip()
                    if text:
                        style_name = "Heading1" if p.name == "h1" else "Heading2" if p.name == "h2" else "Normal"
                        story.append(Paragraph(text, styles[style_name]))
                        story.append(Spacer(1, 4))
                story.append(Spacer(1, 10))
            except Exception as e:
                logger.error(f"Failed to parse epub section {hf}: {e}")
    
    if not story:
        story.append(Paragraph("Empty EPUB file.", styles["Normal"]))
        
    doc.build(story)
    buffer.seek(0)
    return buffer.read()
    
async def epub_to_pdf(contents: bytes) -> bytes:
    return await asyncio.to_thread(_epub_to_pdf_sync, contents)

def _zip_to_pdf_sync(contents: bytes) -> bytes:
    import zipfile
    import fitz
    from PIL import Image
    
    merged_doc = fitz.open()
    
    with zipfile.ZipFile(BytesIO(contents)) as z:
        for filename in sorted(z.namelist()):
            if filename.startswith("__MACOSX") or filename.endswith((".DS_Store", "/")):
                continue
                
            file_bytes = z.read(filename)
            ext = "." + filename.split('.')[-1].lower()
            
            try:
                if ext == ".pdf":
                    doc = fitz.open(stream=file_bytes, filetype="pdf")
                    merged_doc.insert_pdf(doc)
                    doc.close()
                elif ext in (".png", ".jpg", ".jpeg"):
                    doc = fitz.open()
                    img = Image.open(BytesIO(file_bytes))
                    width, height = img.width, img.height
                    page = doc.new_page(width=width, height=height)
                    page.insert_image(fitz.Rect(0, 0, width, height), stream=file_bytes)
                    merged_doc.insert_pdf(doc)
                    doc.close()
                elif ext == ".txt":
                    txt_bytes = _txt_to_pdf_sync(file_bytes)
                    doc = fitz.open(stream=txt_bytes, filetype="pdf")
                    merged_doc.insert_pdf(doc)
                    doc.close()
            except Exception as e:
                logger.error(f"Failed to process zip file {filename}: {e}")
                
    if len(merged_doc) == 0:
        raise ValueError("No convertible files found inside the ZIP archive.")
        
    pdf_bytes = merged_doc.write()
    merged_doc.close()
    return pdf_bytes

async def zip_to_pdf(contents: bytes) -> bytes:
    return await asyncio.to_thread(_zip_to_pdf_sync, contents)


def _pdf_to_zip_sync(contents: bytes) -> bytes:
    import zipfile
    import fitz
    
    doc = fitz.open(stream=contents, filetype="pdf")
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as z:
        for page_num in range(len(doc)):
            page_doc = fitz.open()
            page_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
            page_bytes = page_doc.write()
            page_doc.close()
            z.writestr(f"page_{page_num + 1}.pdf", page_bytes)
    doc.close()
    buffer.seek(0)
    return buffer.read()


async def pdf_to_zip(contents: bytes) -> bytes:
    return await asyncio.to_thread(_pdf_to_zip_sync, contents)


async def convert_office_to_format_via_libreoffice(contents: bytes, src_ext: str, target_ext: str) -> bytes:
    import shutil
    import os
    libreoffice_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice_bin and os.path.exists("/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        libreoffice_bin = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    
    if not libreoffice_bin:
        raise FileNotFoundError("LibreOffice executable not found. Please install LibreOffice.")
        
    with tempfile.NamedTemporaryFile(suffix=f".{src_ext}", delete=False) as tmp_file:
        tmp_file.write(contents)
        input_path = tmp_file.name
        
    output_dir = tempfile.mkdtemp()
    try:
        cmd = [
            libreoffice_bin,
            "--headless",
            "--convert-to",
            target_ext,
            "--outdir",
            output_dir,
            input_path,
        ]
        proc = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice failed: {stderr.decode()}")
        
        result_file = Path(output_dir) / (Path(input_path).stem + f".{target_ext}")
        if not result_file.exists():
            found = list(Path(output_dir).glob(f"*.{target_ext}"))
            if found:
                result_file = found[0]
            else:
                raise FileNotFoundError(f"Converted file .{target_ext} not found in output directory.")
                
        with open(result_file, "rb") as f:
            out_bytes = f.read()
        return out_bytes
    finally:
        Path(input_path).unlink(missing_ok=True)
        shutil.rmtree(output_dir, ignore_errors=True)


async def convert_pdf_to_format_via_libreoffice(contents: bytes, target_ext: str) -> bytes:
    # Route through XLSX/PPTX intermediates for spreadsheet/presentation targets
    if target_ext in ("ods", "xlsx", "xls"):
        xlsx_bytes = await pdf_to_excel(contents)
        return await convert_office_to_format_via_libreoffice(xlsx_bytes, "xlsx", target_ext)
    elif target_ext in ("odp", "pptx", "ppt"):
        pptx_bytes = await pdf_to_ppt(contents)
        return await convert_office_to_format_via_libreoffice(pptx_bytes, "pptx", target_ext)
        
    import shutil
    import os
    libreoffice_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice_bin and os.path.exists("/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        libreoffice_bin = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    
    if not libreoffice_bin:
        raise FileNotFoundError("LibreOffice executable not found. Please install LibreOffice.")
        
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_file:
        tmp_file.write(contents)
        input_path = tmp_file.name
        
    output_dir = tempfile.mkdtemp()
    try:
        infilter = "writer_pdf_import"
        cmd = [
            libreoffice_bin,
            "--headless",
            f"--infilter={infilter}",
            "--convert-to",
            target_ext,
            "--outdir",
            output_dir,
            input_path,
        ]
        proc = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
        
        file_generated = (Path(output_dir) / (Path(input_path).stem + f".{target_ext}")).exists() or bool(list(Path(output_dir).glob(f"*.{target_ext}")))
        
        if proc.returncode != 0 or not file_generated:
            logger.warning(f"LibreOffice primary conversion failed or did not produce output (code={proc.returncode}). Running fallback without infilter.")
            cmd_fallback = [
                libreoffice_bin,
                "--headless",
                "--convert-to",
                target_ext,
                "--outdir",
                output_dir,
                input_path,
            ]
            proc_fallback = await asyncio.create_subprocess_exec(
                *cmd_fallback, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await asyncio.wait_for(proc_fallback.communicate(), timeout=30)
            if proc_fallback.returncode != 0:
                raise RuntimeError(f"LibreOffice fallback failed: {stderr.decode()}")
        
        result_file = Path(output_dir) / (Path(input_path).stem + f".{target_ext}")
        if not result_file.exists():
            found = list(Path(output_dir).glob(f"*.{target_ext}"))
            if found:
                result_file = found[0]
            else:
                raise FileNotFoundError(f"Converted file .{target_ext} not found in output directory.")
                
        with open(result_file, "rb") as f:
            out_bytes = f.read()
        return out_bytes
    finally:
        Path(input_path).unlink(missing_ok=True)
        shutil.rmtree(output_dir, ignore_errors=True)


async def convert_docx_to_format_via_libreoffice(contents: bytes, target_ext: str) -> bytes:
    return await convert_office_to_format_via_libreoffice(contents, "docx", target_ext)


